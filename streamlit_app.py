import streamlit as st
import streamlit_authenticator as stauth
import yaml
from yaml.loader import SafeLoader
import os
import pdfplumber
import pandas as pd
import matplotlib.pyplot as plt
from datetime import datetime
from pathlib import Path
from pptx import Presentation
import openai
import json

# --- Load users from YAML ---
with open('users.yaml') as file:
    config = yaml.load(file, Loader=SafeLoader)

authenticator = stauth.Authenticate(
    config['credentials'],
    config['cookie']['name'],
    config['cookie']['key'],
    config['cookie']['expiry_days']
)

# --- LOGIN ---
login_result = authenticator.login(location='main')
if login_result is None:
    st.stop()

if isinstance(login_result, tuple):
    if len(login_result) == 2:
        name, authentication_status = login_result
        username = getattr(authenticator, "username", None)
    elif len(login_result) == 3:
        name, authentication_status, username = login_result
    else:
        st.error("Unknown login return value structure.")
        st.stop()
else:
    st.error("Unexpected login return type.")
    st.stop()

if authentication_status is False:
    st.error("Username/password is incorrect")
    st.stop()
elif authentication_status is None:
    st.warning("Please enter your username and password")
    st.stop()

# Get user role for instructor dashboard
user_role = config['credentials']['usernames'][username]['role']

st.sidebar.write(f"Logged in as: **{name}** ({username})")
st.success("Login success! You now see the main app.")

# --- Main App ---
st.title("📚 PTA Tutor Chatbot with Quiz & Performance Tracker")

# --- Dynamic course discovery
COURSE_MATERIALS_ROOT = "course_materials"
courses = [d for d in os.listdir(COURSE_MATERIALS_ROOT)
           if os.path.isdir(os.path.join(COURSE_MATERIALS_ROOT, d))]
if not courses:
    courses = ["PTA_1010"]  # fallback

course = st.selectbox("Select your course:", sorted(courses))
course_folder = os.path.join(COURSE_MATERIALS_ROOT, course)

def load_pdf_text(folder):
    text = ""
    if os.path.exists(folder):
        for filename in os.listdir(folder):
            if filename.endswith(".pdf"):
                file_path = os.path.join(folder, filename)
                with pdfplumber.open(file_path) as pdf:
                    for page in pdf.pages:
                        page_text = page.extract_text()
                        if page_text:
                            text += page_text
    return text

pdf_text = load_pdf_text(course_folder)[:3000]

def load_txt_content(folder):
    txt_file = None
    if os.path.exists(folder):
        for filename in os.listdir(folder):
            if filename.endswith(".txt"):
                txt_file = os.path.join(folder, filename)
                break
    if txt_file:
        with open(txt_file, "r", encoding="utf-8") as f:
            return f.read()
    return ""

txt_text = load_txt_content(course_folder)[:3000]

def extract_notes_from_uploaded_pptx(uploaded_file):
    prs = Presentation(uploaded_file)
    all_notes = []
    for i, slide in enumerate(prs.slides):
        slide_title = slide.shapes.title.text if slide.shapes.title else f"Slide {i+1}"
        notes_text = ""
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
        all_notes.append(f"{slide_title}:\n{notes_text}\n")
    return "\n".join(all_notes)

st.sidebar.header("Optional: Upload PowerPoint for Chatbot Content")
uploaded_pptx = st.sidebar.file_uploader("Upload a PowerPoint (.pptx)", type="pptx")
pptx_text = ""
if uploaded_pptx:
    pptx_text = extract_notes_from_uploaded_pptx(uploaded_pptx)
    st.sidebar.success("PowerPoint notes extracted. Chatbot will use these as course content.")

# --- OpenAI setup (NO OpenAI object needed) ---
openai_api_key = st.secrets["openai"]["api_key"]
openai.api_key = openai_api_key

log_path = Path("grading_log.csv")
if not log_path.exists():
    pd.DataFrame(columns=[
        "username", "question_id", "question_text", "user_answer",
        "correct_answer", "correct", "topic", "blooms_level", "confidence", "timestamp"
    ]).to_csv(log_path, index=False)

st.header("💬 Chat with the Tutor")

if "messages" not in st.session_state:
    st.session_state.messages = []

for message in st.session_state.messages:
    with st.chat_message(message["role"]):
        st.markdown(message["content"])

if prompt := st.chat_input("Ask a question about your course..."):
    st.session_state.messages.append({"role": "user", "content": prompt})
    with st.chat_message("user"):
        st.markdown(prompt)

    if pptx_text:
        course_content = pptx_text
        content_source = "PowerPoint notes"
    elif txt_text:
        course_content = txt_text
        content_source = "Text file"
    else:
        course_content = pdf_text
        content_source = "PDF"

    st.info(f"Chatbot is using: {content_source}")

    system_prompt = {
        "role": "system",
        "content": f"""You are a knowledgeable and focused PTA tutor.

Use ONLY this course content to answer questions.

Do NOT reference slide numbers or slide locations in any answers or questions. Only focus on the content itself, not where it appears in the slides or documents.

{course_content}

If the question is unrelated to the material, respond: 'I'm sorry, I can only help with the course content provided.'"""
    }

    try:
        response = openai.chat.completions.create(
            model="gpt-3.5-turbo",
            messages=[system_prompt] + st.session_state.messages
        )
        reply = response.choices[0].message.content
        with st.chat_message("assistant"):
            st.markdown(reply)
        st.session_state.messages.append({"role": "assistant", "content": reply})
    except Exception as e:
        st.error(f"❌ Error: {str(e)}")

# --- Quiz Generator with Blooms Levels 1–5 + Confidence Slider ---
st.header("📝 Quiz Generator")

bloom_option = st.selectbox(
    "Choose Bloom's Taxonomy Level for Quiz:",
    [
        "1 (Recall/Knowledge)",
        "2 (Comprehension)",
        "3 (Application)",
        "4 (Analysis)",
        "5 (Synthesis/Evaluation)",
        "Mixed (Levels 1–5)"
    ]
)

if st.button("Generate Quiz"):
    if pptx_text:
        course_content = pptx_text
    elif txt_text:
        course_content = txt_text
    else:
        course_content = pdf_text

    blooms_level_map = {
        "1": "Recall/Knowledge",
        "2": "Comprehension",
        "3": "Application",
        "4": "Analysis",
        "5": "Synthesis/Evaluation"
    }

    if bloom_option.startswith("Mixed"):
        blooms_instruction = (
            "Generate 5 NPTE-style multiple-choice questions: "
            "one each at Bloom's Level 1 (Recall/Knowledge), "
            "Level 2 (Comprehension), Level 3 (Application), "
            "Level 4 (Analysis), and Level 5 (Synthesis/Evaluation). "
        )
    else:
        level = bloom_option[0]
        level_name = blooms_level_map.get(level, "")
        blooms_instruction = (
            f"Generate 5 NPTE-style multiple-choice questions at Bloom's Level {level} ({level_name}). "
        )

    # Use JSON for parsing (recommended)
    quiz_prompt = (
        f"You are a Physical Therapist Assistant tutor. Based on the following course content, "
        f"{blooms_instruction}"
        "Do NOT reference slide numbers or slide locations in any questions. Focus only on content."
        "For each question, output a JSON object in the following format:\n"
        "{{"
        '"id": "Q001",'
        '"question": "...",'
        '"options": ["A. ...", "B. ...", "C. ...", "D. ..."],'
        '"correct": "A",'
        '"explanation": "...",'
        '"topic": "...",'
        '"blooms_level": "1"'
        "}}\n"
        "Output a list of 5 such questions as a JSON array. Only output JSON, nothing else."
        "\n\n"
        + course_content
    )

    try:
        response = openai.chat.completions.create(
            model="gpt-3.5-turbo",
            messages=[{"role": "user", "content": quiz_prompt}]
        )
        quiz_json = response.choices[0].message.content

        try:
            questions = json.loads(quiz_json)
        except Exception:
            st.error("Could not parse generated quiz! Showing raw output.")
            st.markdown(quiz_json)
            questions = []

        if questions:
            st.session_state['quiz_questions'] = questions
            st.session_state['quiz_idx'] = 0
            st.session_state['quiz_answers'] = []
            st.success("Quiz ready! Click below to start.")
        else:
            st.session_state['quiz_questions'] = []
            st.session_state['quiz_idx'] = 0

    except Exception as e:
        st.error(f"❌ Failed to generate quiz: {str(e)}")

# --- Quiz Taking Flow: One question at a time with confidence tracking ---
if "quiz_questions" in st.session_state and st.session_state["quiz_questions"]:
    idx = st.session_state.get("quiz_idx", 0)
    questions = st.session_state["quiz_questions"]

    if idx < len(questions):
        q = questions[idx]
        st.subheader(f"Question {idx+1}:")
        st.markdown(q["question"])
        user_answer = st.radio("Choose your answer:", q["options"], key=f"q_ans_{idx}")
        confidence = st.slider("How confident are you in your answer?", 1, 5, 3, key=f"conf_{idx}")

        if st.button("Submit Answer", key=f"submit_{idx}"):
            # Record answer
            log_entry = {
                "username": username,
                "question_id": q["id"],
                "question_text": q["question"],
                "user_answer": user_answer[0],  # "A. ...", just take "A"
                "correct_answer": q["correct"],
                "correct": int(user_answer[0] == q["correct"]),
                "topic": q["topic"],
                "blooms_level": q["blooms_level"],
                "confidence": confidence,
                "timestamp": datetime.now().isoformat()
            }
            # Append to grading log
            df = pd.read_csv(log_path)
            df = pd.concat([df, pd.DataFrame([log_entry])], ignore_index=True)
            df.to_csv(log_path, index=False)
            # Save in session_state for feedback display
            st.session_state['quiz_answers'].append({
                "your_answer": user_answer[0],
                "correct": q["correct"],
                "explanation": q.get("explanation", ""),
                "was_correct": user_answer[0] == q["correct"],
                "confidence": confidence,
                "question": q["question"]
            })
            # Go to next question
            st.session_state["quiz_idx"] += 1
            st.experimental_rerun()

    else:
        st.success("Quiz complete! See your answers and explanations below:")
        for i, ans in enumerate(st.session_state['quiz_answers']):
            st.markdown(f"**Question {i+1}:** {ans['question']}")
            st.markdown(f"- Your answer: {ans['your_answer']}")
            st.markdown(f"- Correct answer: {ans['correct']}")
            st.markdown(f"- Correct? {'✅' if ans['was_correct'] else '❌'}")
            st.markdown(f"- Confidence: {ans['confidence']}/5")
            if ans["explanation"]:
                st.markdown(f"**Explanation:** {ans['explanation']}")
            st.markdown("---")
        # Clear quiz after review
        if st.button("Finish Quiz & Clear"):
            del st.session_state["quiz_questions"]
            del st.session_state["quiz_answers"]
            del st.session_state["quiz_idx"]
            st.experimental_rerun()

# --- Instructor Dashboard (admin only) ---
if user_role == "admin":
    st.header("📊 Instructor Dashboard: Flagged/Struggling Students")
    st.markdown("Use controls below to adjust the criteria for flagging students.")

    min_score_frac = st.slider("Minimum score for NOT being flagged (as a fraction)", 0.0, 1.0, 0.7, 0.05)
    min_attempts = st.number_input("Minimum number of quiz attempts to be considered", 1, 20, 3, 1)

    try:
        df = pd.read_csv(log_path)
        if df.empty or "username" not in df.columns or df["username"].isnull().all():
            st.info("No student quiz data yet. Flagged students and analytics will appear after quizzes are taken.")
        else:
            summary = (
                df.groupby("username")["correct"]
                .agg(['mean', 'count'])
                .reset_index()
                .rename(columns={'mean': 'score', 'count': 'attempts'})
            )
            flagged = summary[(summary['score'] < min_score_frac) & (summary['attempts'] >= min_attempts)]
            st.markdown("### 🚩 Students Flagged as Struggling")
            st.dataframe(flagged, use_container_width=True)

            drill_user = st.selectbox("Drill down: select a student", [""] + list(flagged['username']))
            if drill_user:
                st.markdown(f"#### Details for `{drill_user}`")
                user_rows = df[df['username'] == drill_user]
                if 'topic' in user_rows.columns and not user_rows['topic'].isnull().all():
                    st.markdown("**By Topic:**")
                    st.dataframe(user_rows.groupby("topic")["correct"].mean().reset_index().rename(columns={"correct": "Score"}))
                if 'blooms_level' in user_rows.columns and not user_rows['blooms_level'].isnull().all():
                    st.markdown("**By Bloom's Level:**")
                    st.dataframe(user_rows.groupby("blooms_level")["correct"].mean().reset_index().rename(columns={"correct": "Score"}))
                st.markdown("**Raw Results:**")
                st.dataframe(user_rows, use_container_width=True)

            csv = flagged.to_csv(index=False)
            st.download_button("Download flagged report CSV", csv, "flagged_students.csv", "text/csv")
    except Exception as e:
        st.error(f"Admin dashboard error: {e}")

# --- Student Performance Summary (with confidence) ---
with st.expander("📊 Show My Performance Summary", expanded=False):
    try:
        df = pd.read_csv(log_path)
        user_df = df[df["username"] == username] if user_role != "admin" else df
        correct_total = user_df["correct"].sum()
        incorrect_total = len(user_df) - correct_total

        st.write(f"Total Questions Answered: {len(user_df)}")
        st.write(f"✅ Correct: {correct_total}")
        st.write(f"❌ Incorrect: {incorrect_total}")

        if not user_df.empty:
            fig, ax = plt.subplots()
            ax.bar(["Correct", "Incorrect"], [correct_total, incorrect_total])
            ax.set_ylabel("Number of Responses")
            ax.set_title("Student Performance")
            st.pyplot(fig)

            if "confidence" in user_df.columns:
                st.write("**Accuracy by Confidence Level:**")
                conf_summary = user_df.groupby("confidence")["correct"].mean()
                st.bar_chart(conf_summary)

    except Exception as e:
        st.warning("⚠️ No grading data available or error reading log.")
        st.text(str(e))
